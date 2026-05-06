"""
Generate DOCX, PPTX, and PDF files for the Confidence Detector project.
Reads MD source files and produces formatted documents.

Usage: python generate_docs.py
"""

import os
import re
import sys
from pathlib import Path

from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_ORIENT
from docx.oxml.ns import qn, nsdecls
from docx.oxml import parse_xml

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN

BASE_DIR = Path(r"d:\AI Learning\usman\07_Confidence_Detector_Project")


# ─────────────────────────────────────────────
# DOCX GENERATION
# ─────────────────────────────────────────────

def parse_md_lines(md_text):
    """Parse markdown text into structured elements."""
    lines = md_text.split('\n')
    elements = []
    i = 0
    while i < len(lines):
        line = lines[i]

        # Code block
        if line.strip().startswith('```'):
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith('```'):
                code_lines.append(lines[i])
                i += 1
            elements.append(('code', '\n'.join(code_lines)))
            i += 1
            continue

        # Headings
        if line.startswith('#### '):
            elements.append(('h4', line[5:].strip()))
        elif line.startswith('### '):
            elements.append(('h3', line[4:].strip()))
        elif line.startswith('## '):
            elements.append(('h2', line[3:].strip()))
        elif line.startswith('# '):
            elements.append(('h1', line[2:].strip()))

        # Horizontal rule
        elif line.strip() in ('---', '***', '___'):
            elements.append(('hr', ''))

        # Table row
        elif '|' in line and line.strip().startswith('|'):
            # Collect entire table
            table_lines = []
            while i < len(lines) and '|' in lines[i] and lines[i].strip().startswith('|'):
                stripped = lines[i].strip()
                # Skip separator rows like |:---|:---|
                if not re.match(r'^\|[\s:\-|]+\|$', stripped):
                    cells = [c.strip() for c in stripped.split('|')]
                    cells = [c for c in cells if c != '']  # remove empty from leading/trailing |
                    table_lines.append(cells)
                i += 1
            if table_lines:
                elements.append(('table', table_lines))
            continue

        # Blockquote
        elif line.strip().startswith('> '):
            elements.append(('quote', line.strip()[2:]))

        # Bullet list
        elif re.match(r'^[\s]*[-*]\s', line):
            elements.append(('bullet', line.strip().lstrip('-* ').strip()))

        # Numbered list
        elif re.match(r'^[\s]*\d+\.\s', line):
            text = re.sub(r'^[\s]*\d+\.\s*', '', line).strip()
            elements.append(('numbered', text))

        # Empty line
        elif line.strip() == '':
            elements.append(('empty', ''))

        # Regular paragraph
        else:
            elements.append(('para', line))

        i += 1

    return elements


def add_formatted_text(paragraph, text):
    """Add text with bold/italic/code formatting to a paragraph."""
    # Process bold, italic, inline code
    parts = re.split(r'(\*\*[^*]+\*\*|`[^`]+`|\*[^*]+\*)', text)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            run = paragraph.add_run(part[2:-2])
            run.bold = True
        elif part.startswith('`') and part.endswith('`'):
            run = paragraph.add_run(part[1:-1])
            run.font.name = 'Courier New'
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        elif part.startswith('*') and part.endswith('*') and not part.startswith('**'):
            run = paragraph.add_run(part[1:-1])
            run.italic = True
        else:
            paragraph.add_run(part)


def set_cell_shading(cell, color_hex):
    """Set background shading on a table cell."""
    shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color_hex}"/>')
    cell._tc.get_or_add_tcPr().append(shading)


def create_docx(title, md_files, output_path):
    """Create a DOCX file from one or more markdown files."""
    doc = Document()

    # Set default font
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = Pt(11)

    # Set A4 page size
    for section in doc.sections:
        section.page_width = Cm(21.0)
        section.page_height = Cm(29.7)
        section.top_margin = Cm(2.54)
        section.bottom_margin = Cm(2.54)
        section.left_margin = Cm(2.54)
        section.right_margin = Cm(2.54)

    # Title page
    for _ in range(6):
        doc.add_paragraph('')

    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title_para.add_run(title)
    run.bold = True
    run.font.size = Pt(28)
    run.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run('Presentation Confidence Detector')
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(0x64, 0x64, 0xFF)

    doc.add_paragraph('')
    project_line = doc.add_paragraph()
    project_line.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = project_line.add_run('Project Blueprint')
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    doc.add_page_break()

    # Process each MD file
    for md_file in md_files:
        if not os.path.exists(md_file):
            print(f"  WARNING: File not found: {md_file}")
            continue

        with open(md_file, 'r', encoding='utf-8') as f:
            md_text = f.read()

        elements = parse_md_lines(md_text)

        for elem_type, content in elements:
            if elem_type == 'h1':
                p = doc.add_heading(content, level=1)
            elif elem_type == 'h2':
                p = doc.add_heading(content, level=2)
            elif elem_type == 'h3':
                p = doc.add_heading(content, level=3)
            elif elem_type == 'h4':
                p = doc.add_paragraph()
                run = p.add_run(content)
                run.bold = True
                run.font.size = Pt(12)

            elif elem_type == 'code':
                p = doc.add_paragraph()
                p.style = doc.styles['Normal']
                p.paragraph_format.left_indent = Cm(0.5)
                # Add gray background via XML
                pPr = p._p.get_or_add_pPr()
                shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="F0F0F0"/>')
                pPr.append(shading)
                run = p.add_run(content)
                run.font.name = 'Courier New'
                run.font.size = Pt(9)

            elif elem_type == 'table':
                rows = content
                if not rows:
                    continue
                num_cols = max(len(r) for r in rows)
                table = doc.add_table(rows=len(rows), cols=num_cols)
                table.style = 'Light Grid Accent 1'

                for row_idx, row_data in enumerate(rows):
                    for col_idx, cell_text in enumerate(row_data):
                        if col_idx < num_cols:
                            cell = table.cell(row_idx, col_idx)
                            cell.text = ''
                            p = cell.paragraphs[0]
                            # Strip bold markers for table cells
                            clean_text = re.sub(r'\*\*([^*]+)\*\*', r'\1', cell_text)
                            clean_text = re.sub(r'`([^`]+)`', r'\1', clean_text)
                            add_formatted_text(p, cell_text)
                            p.paragraph_format.space_after = Pt(2)

                            # Header row styling
                            if row_idx == 0:
                                for run in p.runs:
                                    run.bold = True
                                set_cell_shading(cell, "1a1a2e")
                                for run in p.runs:
                                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

                doc.add_paragraph('')  # spacing after table

            elif elem_type == 'quote':
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = Cm(1.0)
                p.paragraph_format.space_before = Pt(6)
                p.paragraph_format.space_after = Pt(6)
                pPr = p._p.get_or_add_pPr()
                shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="E8E8F0"/>')
                pPr.append(shading)
                add_formatted_text(p, content)
                for run in p.runs:
                    run.italic = True

            elif elem_type == 'bullet':
                p = doc.add_paragraph(style='List Bullet')
                p.text = ''
                add_formatted_text(p, content)

            elif elem_type == 'numbered':
                p = doc.add_paragraph(style='List Number')
                p.text = ''
                add_formatted_text(p, content)

            elif elem_type == 'para':
                if content.strip():
                    p = doc.add_paragraph()
                    add_formatted_text(p, content)

            elif elem_type == 'hr':
                # Add a thin line as separator
                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(6)
                p.paragraph_format.space_after = Pt(6)
                run = p.add_run('_' * 80)
                run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
                run.font.size = Pt(6)

            elif elem_type == 'empty':
                pass  # skip empty lines to avoid excessive spacing

    doc.save(str(output_path))
    print(f"  Created: {output_path}")


# ─────────────────────────────────────────────
# PPTX GENERATION
# ─────────────────────────────────────────────

BG_COLOR = PptxRGBColor(20, 20, 30)
WHITE = PptxRGBColor(255, 255, 255)
BLUE_ACCENT = PptxRGBColor(100, 100, 255)
LIGHT_GRAY = PptxRGBColor(180, 180, 190)
DARK_GRAY = PptxRGBColor(40, 40, 55)


def set_slide_bg(slide, color):
    """Set the background color of a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(
        PptxInches(left), PptxInches(top),
        PptxInches(width), PptxInches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = PptxPt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def add_bullet_slide(prs, title_text, bullets, subtitle=None):
    """Add a slide with a title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_slide_bg(slide, BG_COLOR)

    # Title
    add_text_box(slide, 0.5, 0.3, 9, 0.8, title_text,
                 font_size=28, bold=True, color=BLUE_ACCENT)

    # Optional subtitle
    y_start = 1.2
    if subtitle:
        add_text_box(slide, 0.5, 1.1, 9, 0.5, subtitle,
                     font_size=14, color=LIGHT_GRAY)
        y_start = 1.6

    # Bullets
    txBox = slide.shapes.add_textbox(
        PptxInches(0.7), PptxInches(y_start),
        PptxInches(8.6), PptxInches(5.5 - (y_start - 1.2))
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.space_after = PptxPt(8)

        # Check if bullet has a bold prefix (e.g., "Title: description")
        if ':' in bullet and not bullet.startswith('http'):
            parts = bullet.split(':', 1)
            run_bold = p.add_run()
            run_bold.text = parts[0] + ':'
            run_bold.font.size = PptxPt(18)
            run_bold.font.bold = True
            run_bold.font.color.rgb = WHITE

            run_rest = p.add_run()
            run_rest.text = parts[1]
            run_rest.font.size = PptxPt(18)
            run_rest.font.color.rgb = LIGHT_GRAY
        else:
            run = p.add_run()
            run.text = bullet
            run.font.size = PptxPt(18)
            run.font.color.rgb = WHITE

    return slide


def add_table_slide(prs, title_text, headers, rows):
    """Add a slide with a title and table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    add_text_box(slide, 0.5, 0.3, 9, 0.8, title_text,
                 font_size=28, bold=True, color=BLUE_ACCENT)

    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        PptxInches(0.4), PptxInches(1.3),
        PptxInches(9.2), PptxInches(5.0)
    )
    table = table_shape.table

    # Style header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = PptxPt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
        # Dark blue header background
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = PptxRGBColor(30, 30, 80)

    # Data rows
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = PptxPt(12)
                paragraph.font.color.rgb = LIGHT_GRAY
            cell_fill = cell.fill
            cell_fill.solid()
            if i % 2 == 0:
                cell_fill.fore_color.rgb = PptxRGBColor(30, 30, 45)
            else:
                cell_fill.fore_color.rgb = DARK_GRAY

    return slide


def create_master_pptx(output_path):
    """Create the master PPTX presentation."""
    prs = Presentation()
    prs.slide_width = PptxInches(10)
    prs.slide_height = PptxInches(7.5)

    # ── Slide 1: Title ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)
    add_text_box(slide, 1, 2.0, 8, 1.2,
                 "Presentation Confidence Detector",
                 font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 3.3, 8, 0.8,
                 "Project Blueprint",
                 font_size=24, color=BLUE_ACCENT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.5, 8, 0.6,
                 "AI-powered real-time presentation coaching system",
                 font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ── Slide 2: The Problem ──
    add_bullet_slide(prs, "The Problem", [
        "Professionals practice alone with zero objective feedback",
        "People are terrible at judging their own confidence",
        "Hiring a presentation coach costs $200-500/hour",
        "Recording yourself takes 2x the time and you still miss things",
        "Without data, you are guessing how you performed",
        "Result: unprepared speakers, lost deals, missed opportunities"
    ])

    # ── Slide 3: How It Works ──
    add_bullet_slide(prs, "How It Works", [
        "1. User starts a practice session (topic + start button)",
        "2. Camera captures face at 30 FPS for expression analysis",
        "3. Microphone captures voice for speech + audio analysis",
        "4. Four engines run simultaneously: Face, STT, NLP, Voice",
        "5. Scoring algorithm combines signals into one score (0-100)",
        "6. Live dashboard shows score + coaching nudges in real-time",
        "7. Session report gives detailed breakdown after stopping"
    ], subtitle="Pipeline: Camera + Mic -> 4 Engines -> Score -> Feedback")

    # ── Slide 4: Architecture ──
    add_bullet_slide(prs, "System Architecture (5 Layers)", [
        "Layer 1 - Input: Camera (30 FPS) + Microphone (44.1 kHz)",
        "Layer 2 - Processing: Face, STT, NLP, Voice engines in parallel",
        "Layer 3 - Scoring: Weighted combination -> single 0-100 score",
        "Layer 4 - Feedback: Live dashboard + coaching alerts",
        "Layer 5 - Report: Post-session summary with charts and advice"
    ], subtitle="Each engine runs independently -- if one fails, others keep working")

    # ── Slide 5: Scoring Formula ──
    add_bullet_slide(prs, "Scoring Formula", [
        "Final Score = (Face x 0.40) + (Speech x 0.35) + (Voice x 0.25)",
        "Face (0.40): Eye contact, expression, blink rate, posture",
        "Speech (0.35): Filler words, hedging, pace, repetitions",
        "Voice (0.25): Pitch stability, volume, silence ratio",
        "Smoothed with EMA: 70% stability + 30% new reading",
        "Score Labels: 85-100 Highly Confident, 70-84 Confident, 50-69 Moderate, <50 Developing"
    ], subtitle="Weights based on signal reliability + audience impact")

    # ── Slide 6: Detection Map ──
    add_table_slide(prs, "Detection Map (6 Sub-Problems)",
        ["#", "Sub-Problem", "Technology", "Update Rate"],
        [
            ["1", "Facial Expressions", "MediaPipe FaceMesh (468 pts)", "Every frame"],
            ["2", "Eye Contact", "Iris + head landmarks", "Every frame"],
            ["3", "Filler Words", "STT + pattern matching", "Per sentence"],
            ["4", "Hedging Language", "Phrase list matching", "Per sentence"],
            ["5", "Speaking Pace", "Word count / time", "Every 5 sec"],
            ["6", "Voice Steadiness", "FFT audio analysis", "Continuous"],
        ]
    )

    # ── Slide 7: Technology Choices ──
    add_bullet_slide(prs, "Technology Choices", [
        "MediaPipe FaceMesh: Free, fast, 468 landmarks, runs offline in browser",
        "Web Speech API (MVP): Zero setup, instant prototype, Chrome/Edge",
        "Vosk (Production): Offline, cross-browser, no timeouts, ~200MB/user",
        "Pattern Matching (NLP): Instant, free, 85%+ accuracy for fillers",
        "Web Audio API + FFT: Built into every browser, pitch + volume",
        "Claude API (v2): Deep coaching feedback, multimodal analysis"
    ], subtitle="WHY each choice: Speed + Cost + Offline capability")

    # ── Slide 8: Phase 1 ──
    add_bullet_slide(prs, "Phase 1: Hardware Access", [
        "getUserMedia API: Access camera + microphone in browser",
        "MediaStream Tracks: Separate video/audio for independent routing",
        "Permission Handling: Graceful errors + fallback modes",
        "Stream Cleanup: Prevent camera light staying on (learned from ExamGuard)",
        "Canvas API: Draw face landmarks + feedback overlays on video"
    ], subtitle="Foundation: no hardware access = no data to analyze")

    # ── Slide 9: Phase 2 ──
    add_bullet_slide(prs, "Phase 2: Speech & Audio", [
        "Speech-to-Text: Convert spoken words to analyzable text",
        "Interim vs Final results: Show live text, analyze only confirmed",
        "Auto-restart pattern: Keep STT alive during 5-30 min sessions",
        "Web Audio API + FFT: Extract pitch, volume from raw audio",
        "Volume/pitch tracking over time: Detect trends and tremor",
        "Silence detection: Distinguish thinking pauses from nervous gaps"
    ], subtitle="Entirely NEW -- ExamGuard had zero audio processing")

    # ── Slide 10: Phase 3 ──
    add_bullet_slide(prs, "Phase 3: Face Expression", [
        "MediaPipe FaceMesh: 468 landmarks + 52 blendshapes in browser",
        "Blink rate detection: 15-20/min normal, >25/min stress signal",
        "Eye contact tracking: Iris + head direction, 60%+ target",
        "Expression classification: Neutral, happy, tense, worried",
        "Genuine vs forced smile: Duchenne marker (AU6 + AU12)",
        "Baseline calibration: Personalize to each user's resting face"
    ], subtitle="50% reused from ExamGuard, 50% new skills")

    # ── Slide 11: Phase 4 ──
    add_bullet_slide(prs, "Phase 4: NLP Text Analysis", [
        "Tokenization: Split transcript into clean words",
        "Filler detection: 18 fillers (um, uh, like, basically...)",
        "Hedging detection: 16 phrases (I think, maybe, sort of...)",
        "Repetition detection: Stutters + phrase repeats",
        "Speaking pace: WPM calculation with 30-sec rolling window",
        "Text scoring: Weighted combination -> speech confidence 0-100"
    ], subtitle="100% new territory -- text analysis for confidence signals")

    # ── Slide 12: Phase 5 ──
    add_bullet_slide(prs, "Phase 5: System Integration", [
        "Multi-engine orchestration: 4 engines running in parallel",
        "Scoring algorithm: Weighted combination with EMA smoothing",
        "State management: React useRef for high-frequency data",
        "Real-time dashboard: Gauge, transcript, charts, waveform",
        "Coaching alerts: Rule-based nudges with cooldown timers",
        "Session recording + report generation"
    ], subtitle="Where demos become a product")

    # ── Slide 13: MVP Versions ──
    add_table_slide(prs, "MVP Versions",
        ["Version", "What It Proves", "Engines", "Timeline"],
        [
            ["v0.1", "MediaPipe works in browser", "Face only", "3-5 days"],
            ["v0.2", "Two engines run together", "Face + STT + NLP", "1 week"],
            ["v0.3", "Full detection system works", "All 4 engines", "2-3 weeks"],
            ["v1.0", "Production-quality product", "All + AI + History", "2 weeks"],
        ]
    )

    # ── Slide 14: What's Next ──
    add_bullet_slide(prs, "What's Next", [
        "Build v0.1 (Face Only) -- prove MediaPipe works in browser",
        "Validate camera + canvas pipeline with real face mesh",
        "Then add speech (v0.2), full system (v0.3), polish (v1.0)",
        "Each version is a WORKING demo you can show",
        "When blocked: cut scope, not quality",
        "Goal: a portfolio-quality AI product you built yourself"
    ], subtitle="Same philosophy as ExamGuard: build in layers, ship working demos")

    prs.save(str(output_path))
    print(f"  Created: {output_path}")


# ─────────────────────────────────────────────
# PDF GENERATION
# ─────────────────────────────────────────────

def try_convert_to_pdf(file_paths):
    """Attempt to convert DOCX/PPTX files to PDF."""
    try:
        import docx2pdf
        print("\n  docx2pdf found. Attempting PDF conversion...")
        converted = 0
        for fpath in file_paths:
            fpath = str(fpath)
            if fpath.endswith('.docx') or fpath.endswith('.pptx'):
                pdf_path = fpath.rsplit('.', 1)[0] + '.pdf'
                try:
                    docx2pdf.convert(fpath, pdf_path)
                    print(f"    PDF created: {pdf_path}")
                    converted += 1
                except Exception as e:
                    print(f"    Failed to convert {os.path.basename(fpath)}: {e}")
        if converted > 0:
            print(f"  Successfully converted {converted} file(s) to PDF.")
        else:
            print("  No files were converted to PDF.")
    except ImportError:
        print("\n  docx2pdf not installed. Attempting to install...")
        try:
            import subprocess
            subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'docx2pdf', '-q'])
            import docx2pdf
            print("  docx2pdf installed. Attempting PDF conversion...")
            converted = 0
            for fpath in file_paths:
                fpath = str(fpath)
                if fpath.endswith('.docx') or fpath.endswith('.pptx'):
                    pdf_path = fpath.rsplit('.', 1)[0] + '.pdf'
                    try:
                        docx2pdf.convert(fpath, pdf_path)
                        print(f"    PDF created: {pdf_path}")
                        converted += 1
                    except Exception as e:
                        print(f"    Failed to convert {os.path.basename(fpath)}: {e}")
            if converted == 0:
                print("  PDF conversion requires Microsoft Word/PowerPoint installed.")
                print("  To convert manually: Open each file in Word/PowerPoint -> Save As PDF")
        except Exception as e:
            print(f"  Could not install docx2pdf: {e}")
            print("  PDF conversion requires Microsoft Word installed.")
            print("  To convert manually: Open each file in Word/PowerPoint -> Save As PDF")


# ─────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────

def main():
    print("=" * 60)
    print("  Confidence Detector Document Generator")
    print("=" * 60)

    all_generated = []

    # ── DOCX FILES ──
    print("\n--- Generating DOCX files ---\n")

    docx_configs = [
        {
            "title": "Overview",
            "folder": "00_Overview",
            "md_files": [
                "01_What_is_Confidence_Detector.md",
                "02_System_Architecture.md",
                "03_Detection_Map.md",
                "04_Scoring_Explained.md",
            ],
            "output": "00_Overview.docx"
        },
        {
            "title": "Phase 1: Hardware & Media Access",
            "folder": "01_Phase_Hardware_Access",
            "md_files": ["01_Hardware_Media_Access.md"],
            "output": "01_Phase_Hardware_Access.docx"
        },
        {
            "title": "Phase 2: Speech & Audio Processing",
            "folder": "02_Phase_Speech_and_Audio",
            "md_files": ["01_Speech_and_Audio.md"],
            "output": "02_Phase_Speech_and_Audio.docx"
        },
        {
            "title": "Phase 3: Face & Expression Detection",
            "folder": "03_Phase_Face_Expression",
            "md_files": ["01_Face_Expression_Detection.md"],
            "output": "03_Phase_Face_Expression.docx"
        },
        {
            "title": "Phase 4: NLP & Text Analysis",
            "folder": "04_Phase_NLP_Text_Analysis",
            "md_files": ["01_NLP_Text_Analysis.md"],
            "output": "04_Phase_NLP_Text_Analysis.docx"
        },
        {
            "title": "Phase 5: System Integration & UI",
            "folder": "05_Phase_System_Integration",
            "md_files": ["01_System_Integration.md"],
            "output": "05_Phase_System_Integration.docx"
        },
        {
            "title": "MVP Versions",
            "folder": "06_MVP_Versions",
            "md_files": ["01_MVP_Versions.md"],
            "output": "06_MVP_Versions.docx"
        },
    ]

    for cfg in docx_configs:
        folder_path = BASE_DIR / cfg["folder"]
        md_paths = [folder_path / f for f in cfg["md_files"]]
        output_path = folder_path / cfg["output"]
        print(f"  Generating {cfg['output']}...")
        create_docx(cfg["title"], md_paths, output_path)
        all_generated.append(output_path)

    # ── PPTX FILE ──
    print("\n--- Generating Master PPTX ---\n")
    pptx_path = BASE_DIR / "07_Confidence_Detector_Project.pptx"
    print(f"  Generating {pptx_path.name}...")
    create_master_pptx(pptx_path)
    all_generated.append(pptx_path)

    # ── PDF CONVERSION ──
    print("\n--- Attempting PDF Conversion ---")
    try_convert_to_pdf(all_generated)

    # ── SUMMARY ──
    print("\n" + "=" * 60)
    print("  SUMMARY")
    print("=" * 60)
    print(f"\n  DOCX files generated: {len(docx_configs)}")
    print(f"  PPTX files generated: 1")
    print(f"\n  All files:")
    for f in all_generated:
        size_kb = os.path.getsize(f) / 1024 if os.path.exists(f) else 0
        print(f"    {f.relative_to(BASE_DIR)}  ({size_kb:.0f} KB)")

    print("\n  Done!\n")


if __name__ == '__main__':
    main()
