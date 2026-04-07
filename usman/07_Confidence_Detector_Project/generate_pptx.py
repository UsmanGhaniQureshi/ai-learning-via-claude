"""
Generate a comprehensive 24-slide PowerPoint presentation for the
Presentation Confidence Detector project.

Dark theme, widescreen (13.333 x 7.5 inches), standalone document.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ── Colors ──────────────────────────────────────────────────────────
BG       = RGBColor(15, 15, 25)
WHITE    = RGBColor(255, 255, 255)
BLUE     = RGBColor(100, 130, 255)
GREEN    = RGBColor(0, 200, 100)
YELLOW   = RGBColor(255, 210, 0)
RED      = RGBColor(255, 60, 60)
GRAY     = RGBColor(180, 180, 190)
DIM_GRAY = RGBColor(140, 140, 155)

# ── Dimensions ──────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers ─────────────────────────────────────────────────────────

def set_slide_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    return txBox.text_frame


def add_line(tf, text, size=17, color=GRAY, bold=False, italic=False, alignment=PP_ALIGN.LEFT, space_before=0, space_after=0):
    """Add a paragraph to an existing text frame."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p


def title_line(tf, text):
    add_line(tf, text, size=32, color=BLUE, bold=True, space_after=6)


def header_line(tf, text):
    add_line(tf, text, size=22, color=YELLOW, bold=True, space_before=4, space_after=2)


def body_line(tf, text):
    add_line(tf, text, size=17, color=GRAY, space_after=2)


def warn_line(tf, text):
    add_line(tf, text, size=16, color=RED, bold=True, space_before=2, space_after=2)


def note_line(tf, text):
    add_line(tf, text, size=15, color=DIM_GRAY, italic=True, space_before=2, space_after=2)


def green_line(tf, text, size=17):
    add_line(tf, text, size=size, color=GREEN, bold=True, space_after=2)


# ── Slide builder ───────────────────────────────────────────────────

def make_slide(prs):
    """Return a blank slide with dark background."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide)
    return slide


def content_frame(slide, left=0.6, top=0.5, width=12.1, height=6.6):
    return add_textbox(slide, Inches(left), Inches(top), Inches(width), Inches(height))


# ====================================================================
#  BUILD PRESENTATION
# ====================================================================

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ── Slide 1: Title ──────────────────────────────────────────────────
s = make_slide(prs)
tf = content_frame(s, top=1.8, height=4.5)
add_line(tf, "Presentation Confidence Detector", size=44, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER, space_after=10)
add_line(tf, "Complete Project Blueprint", size=28, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER, space_after=20)
add_line(tf, "A browser-based AI system that watches your face, listens to your speech,", size=18, color=GRAY, alignment=PP_ALIGN.CENTER, space_after=2)
add_line(tf, "and analyzes your voice to give you a real-time confidence score (0-100).", size=18, color=GRAY, alignment=PP_ALIGN.CENTER, space_after=14)
add_line(tf, "Everything runs locally. No server required for core detection.", size=16, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER, space_after=6)
add_line(tf, "24 slides  |  Standalone document  |  Read this and understand the entire project", size=15, color=DIM_GRAY, italic=True, alignment=PP_ALIGN.CENTER)


# ── Slide 2: The Problem ────────────────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "The Problem")
header_line(tf, ">> Why does this project exist?")
body_line(tf, "1. People cannot objectively judge their own confidence while presenting.")
body_line(tf, "   You THINK you look confident. The audience sees something different.")
body_line(tf, "2. Professional presentation coaches cost $200-500 per hour.")
body_line(tf, "   Most people cannot afford 10 sessions to improve.")
body_line(tf, "3. No tool gives objective, real-time feedback during practice.")
body_line(tf, "   Recording yourself and watching later is slow and subjective.")
header_line(tf, ">> What we build instead")
body_line(tf, "A free, offline, browser-based tool that acts as your AI presentation coach.")
body_line(tf, "It watches your face, listens to your words, analyzes your voice, and gives")
body_line(tf, "you a real-time confidence score with coaching nudges -- all in your browser.")
note_line(tf, "\"If you can't measure it, you can't improve it.\" This tool measures confidence.")


# ── Slide 3: How It Works ───────────────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "How It Works: 7-Step Pipeline")
header_line(tf, ">> From button click to coaching feedback")
body_line(tf, "Step 1: User clicks START  -->  Camera + Microphone activated")
body_line(tf, "Step 2: Camera sends 30 frames/sec  -->  Microphone streams audio")
body_line(tf, "Step 3: Face Engine extracts expressions, eye contact, blinks, posture")
body_line(tf, "Step 4: STT Engine converts speech to text  -->  NLP Engine finds fillers & hedges")
body_line(tf, "Step 5: Voice Engine measures pitch, volume, silence ratio")
body_line(tf, "Step 6: Scoring Algorithm combines all signals into one score (0-100)")
body_line(tf, "Step 7: Dashboard shows live score, coaching alerts, transcript with highlights")
header_line(tf, ">> After the session")
body_line(tf, "User clicks STOP  -->  Session Report generated with timeline, breakdown, top 3 tips")
note_line(tf, "All 7 steps happen continuously. The user just sees one smooth experience.")


# ── Slide 4: Architecture ───────────────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "System Architecture: 5 Layers")
header_line(tf, ">> Layer 1: Input Layer")
body_line(tf, "Camera (30 FPS video frames) + Microphone (44,100 audio samples/sec)")
header_line(tf, ">> Layer 2: Processing Layer (4 Engines)")
body_line(tf, "Face Engine (MediaPipe)  |  STT Engine (speech-to-text)  |  NLP Engine (text analysis)  |  Voice Engine (FFT)")
header_line(tf, ">> Layer 3: Scoring Layer")
body_line(tf, "Face Score x 0.40  +  Speech Score x 0.35  +  Voice Score x 0.25  =  Confidence (0-100)")
header_line(tf, ">> Layer 4: Feedback Layer")
body_line(tf, "Animated confidence meter  |  Coaching alerts (max 1 per 15s)  |  Live transcript  |  Metrics cards")
header_line(tf, ">> Layer 5: Report Layer")
body_line(tf, "Final score + per-engine breakdown + timeline chart + filler analysis + top 3 improvements")
note_line(tf, "Each engine runs independently. If one fails, the others keep working. Degrade gracefully, never crash.")


# ── Slide 5: The Scoring Formula ────────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "The Scoring Formula")
header_line(tf, ">> The formula")
green_line(tf, "Final Score = (Face x 0.40) + (Speech x 0.35) + (Voice x 0.25)", size=20)
header_line(tf, ">> WHY these weights? Three reasons:")
body_line(tf, "1. MEHRABIAN'S RESEARCH: ~55% of perceived confidence is visual (face/body).")
body_line(tf, "   Face gets 0.40 (discounted from 0.55 because we only see the face, not full body).")
body_line(tf, "2. RELIABILITY: Face detection (MediaPipe) is the most accurate engine we have.")
body_line(tf, "   Voice analysis in a browser is the least reliable (mic quality varies). So voice = 0.25.")
body_line(tf, "3. AUDIENCE IMPACT: What you SAY is what the audience remembers after the talk.")
body_line(tf, "   Speech (fillers, hedges, pace) is concrete and countable. Speech = 0.35.")
warn_line(tf, "IMPORTANT: These are starting heuristics, NOT scientific constants.")
note_line(tf, "After testing with real users, adjust weights based on data. Equal weights (0.33 each) assume equal reliability -- they don't have it.")


# ── Slide 6: Face Score Deep Dive ───────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "Face Score Deep Dive")
header_line(tf, ">> Baseline: 45 points. Then adjust up or down:")
body_line(tf, "EXPRESSION:    Smiling +25  |  Speaking +15  |  Neutral 0  |  Worried -10  |  Tense -20")
body_line(tf, "EYE CONTACT:   >80% of time looking at camera: +20  |  60-80%: +12  |  30-60%: +5  |  <15%: -15")
body_line(tf, "BLINK RATE:    <25/min: 0  |  25-35/min: -5  |  >35/min: -10  (normal is 15-20/min)")
body_line(tf, "POSTURE:       Upright: +8  |  Tilted: -5  |  Slouching: -12  (shoulder landmarks)")
body_line(tf, "FIDGETING:     Score <25: 0  |  25-50: -5  |  >50: -10  (shoulder/wrist movement rate)")
body_line(tf, "HAND GESTURES: Raised/gesturing: +7  (positive movement vs nervous fidgeting)")
header_line(tf, ">> Example calculation")
body_line(tf, "45 (base) + 15 (speaking) + 12 (75% eye contact) + 0 (normal blinks) + 8 (upright) + 0 (no fidget) + 3 (some gestures) = 83")
note_line(tf, "Same body part, different meaning: hands gesturing = +7, hands fidgeting = -10. HOW it moves matters.")


# ── Slide 7: WHY Blendshapes ────────────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "WHY Blendshapes (Our Real Experience)")
header_line(tf, ">> What we tried first: Manual distance calculations")
body_line(tf, "We measured pixel distances between facial landmarks (mouth corners, eye edges).")
body_line(tf, "Result: FAILED BADLY.")
warn_line(tf, "- Our smile detection gave NEGATIVE scores when people smiled")
warn_line(tf, "  (mouth opening shifted the center point and broke the distance math)")
warn_line(tf, "- 'Tense' triggered on every normal resting face (lip distance threshold too tight)")
warn_line(tf, "- Thresholds that worked on one person failed on another")
header_line(tf, ">> What works: MediaPipe Blendshapes")
body_line(tf, "MediaPipe provides 52 pre-computed blendshape scores (0.0 to 1.0) for each face.")
body_line(tf, "Trained on millions of faces. Already knows what a smile looks like across different")
body_line(tf, "face shapes, sizes, and lighting. We just read the score: mouthSmileLeft > 0.08 = smiling.")
green_line(tf, "Blendshapes work because the hard problem (what is a smile?) is already solved.")
note_line(tf, "We also considered FER (emotion classifier) but research shows 40-60% accuracy on subtle expressions. Too unreliable.")


# ── Slide 8: Speech Score Deep Dive ─────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "Speech Score Deep Dive")
header_line(tf, ">> Start at 100. Subtract penalties:")
green_line(tf, "Speech Score = 100 - FillerPenalty - HedgePenalty - RepetitionPenalty - PacePenalty")
body_line(tf, "FILLERS:      penalty = filler_rate% x 5, capped at -40")
body_line(tf, "              Example: 6 fillers in 280 words = 2.1% rate --> penalty = 10")
body_line(tf, "HEDGES:       penalty = count x 3, capped at -30")
body_line(tf, "              Each 'I think', 'maybe', 'sort of' costs 3 points")
body_line(tf, "REPETITIONS:  penalty = count x 5, capped at -15")
body_line(tf, "              'the the main point' = lost train of thought")
body_line(tf, "PACE:         130-160 WPM = 0 penalty | 100-130 or 160-180 = -5 | <100 or >180 = -15")
header_line(tf, ">> WHY caps on each penalty?")
body_line(tf, "If fillers alone could take score to 0, the other signals would be meaningless.")
note_line(tf, "Example: 100 - 10 (fillers) - 6 (2 hedges) - 0 (no repeats) - 0 (140 WPM) = 84")


# ── Slide 9: The Filler Word Challenge ──────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "The Filler Word Challenge")
header_line(tf, ">> The core problem: context")
body_line(tf, "'I like pizza'        -->  'like' is a REAL WORD  (should not penalize)")
body_line(tf, "'It was, like, good'  -->  'like' is a FILLER    (should penalize)")
body_line(tf, "Pattern matching cannot tell the difference. It sees 'like' and counts it.")
header_line(tf, ">> Our decision: Accept ~15% false positive rate")
body_line(tf, "In presentations, 'like' as filler is FAR more common than 'like' as verb.")
body_line(tf, "Catching 85% of real fillers instantly and for free is better than")
body_line(tf, "catching 98% with a 3-second delay and API cost.")
header_line(tf, ">> v2 solution: Claude API for context")
body_line(tf, "Send the sentence to Claude: 'Is like a filler or real word here?'")
body_line(tf, "98% accuracy, but costs money and adds latency. Optional toggle in v2.")
warn_line(tf, "For v1: pattern matching is instant, free, and 85% accurate. Ship it.")
note_line(tf, "Same trade-off applies to 'so', 'basically', 'actually' -- sometimes real, usually fillers in presentations.")


# ── Slide 10: Voice Score Deep Dive ─────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "Voice Score Deep Dive")
header_line(tf, ">> Three signals from raw audio:")
body_line(tf, "1. VOLUME CONSISTENCY: Is the voice steady or fading at end of sentences?")
body_line(tf, "   >40% drop from average = trailing off (low confidence signal)")
body_line(tf, "2. PITCH VARIATION (FFT): Monotone (StdDev <10 Hz) vs natural (20-50 Hz) vs shaky (>80 Hz)")
body_line(tf, "   Confident speakers vary pitch naturally. Nervous speakers have rapid pitch swings.")
body_line(tf, "3. SILENCE RATIO: 1-3 sec pauses = thinking (OK). >5 sec = lost for words (bad).")
body_line(tf, "   15% silence ratio is normal. Above 30% signals struggle.")
header_line(tf, ">> WHY weighted lowest (0.25)?")
body_line(tf, "Browser audio analysis has real limitations:")
body_line(tf, "- Mic quality varies wildly between devices")
body_line(tf, "- Background noise corrupts readings")
body_line(tf, "- Distance from mic changes volume independent of confidence")
body_line(tf, "- Browser FFT has lower precision than dedicated audio tools")
warn_line(tf, "Voice is a SUPPORTING signal. If face + speech say 'confident' but voice says 'shaky', score drops slightly -- not dramatically.")


# ── Slide 11: Scoring Walkthrough ───────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "Scoring Walkthrough: Real Numbers")
header_line(tf, ">> Scenario: 2-minute practice session")
body_line(tf, "FACE: Speaking 60%, neutral 30%, smiling 10% | 75% eye contact | 22 blinks/min | Upright | Minimal fidget")
body_line(tf, "  --> 45 (base) + 15 (speaking) + 12 (eye) + 0 (blinks) + 8 (posture) + 0 (fidget) + 3 (gestures) = 83")
body_line(tf, "SPEECH: 280 words, 140 WPM | 6 fillers (2.1%) | 2 hedges | 0 repetitions")
body_line(tf, "  --> 100 - 10 (fillers: 2.1x5) - 6 (hedges: 2x3) - 0 (repeats) - 0 (pace) = 84")
body_line(tf, "VOICE: Steady volume | Natural pitch variation | 15% silence ratio")
body_line(tf, "  --> Volume 85 + Pitch 80 + Silence 90 = average 85")
header_line(tf, ">> Combined Final Score")
green_line(tf, "Final = (83 x 0.40) + (84 x 0.35) + (85 x 0.25) = 33.2 + 29.4 + 21.25 = 83.85 ~ 84", size=19)
green_line(tf, "Verdict: 84 = 'Confident' -- good eye contact, minimal fillers, steady voice")
note_line(tf, "Room to improve: reduce the 6 fillers and 2 hedges to push above 90.")


# ── Slide 12: Score Smoothing (EMA) ─────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "Score Smoothing: Exponential Moving Average (EMA)")
header_line(tf, ">> The problem without smoothing")
body_line(tf, "Raw scores jump every frame: 72 -> 45 -> 81 -> 50. Looks broken. Useless for the user.")
header_line(tf, ">> Analogy: Weather forecast")
body_line(tf, "A weather app doesn't say '25C! Now 23C! Now 26C!' every minute.")
body_line(tf, "It smooths: 'Today around 24-25C.' We do the same with confidence scores.")
header_line(tf, ">> The formula")
green_line(tf, "Displayed Score = (0.3 x latest reading) + (0.7 x previous displayed score)")
body_line(tf, "70% stability (where it WAS) + 30% responsiveness (new reading)")
header_line(tf, ">> Worked example")
body_line(tf, "Frame 1: Raw=75  --> Display: 75 (first reading)")
body_line(tf, "Frame 2: Raw=60  --> Display: (0.3x60) + (0.7x75) = 18 + 52.5 = 71")
body_line(tf, "Frame 3: Raw=80  --> Display: (0.3x80) + (0.7x71) = 24 + 49.7 = 74")
body_line(tf, "Frame 4: Raw=50  --> Display: (0.3x50) + (0.7x74) = 15 + 51.8 = 67")
note_line(tf, "Instead of jumping 75->60->80->50, displayed score moves smoothly 75->71->74->67.")


# ── Slide 13: Technology Choices ────────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "Technology Choices (and WHY)")
header_line(tf, ">> WHY MediaPipe Face Mesh?")
body_line(tf, "Free, fast (30+ FPS), 468 landmarks + 52 blendshapes, runs offline, proven in ExamGuard")
body_line(tf, "WHY NOT FER/AffectNet: 40-60% accuracy on subtle expressions. Black box, can't tune.")
header_line(tf, ">> WHY Vosk (offline STT)?")
body_line(tf, "Runs 100% offline, no API cost, supports multiple languages, scalable to many users")
body_line(tf, "WHY NOT paid APIs (Google/AWS): Cost per minute, requires internet, privacy concerns")
body_line(tf, "FALLBACK: Web Speech API for Chrome/Edge MVP (simpler setup, Chrome-only)")
header_line(tf, ">> WHY Pattern Matching for NLP?")
body_line(tf, "Instant execution (no latency), free (no API), 85% accuracy on fillers/hedges")
body_line(tf, "WHY NOT Claude API for every sentence: 2-5 sec latency per call, costs money")
body_line(tf, "v2 PLAN: Add Claude as optional layer for context-aware analysis (98% accuracy)")
header_line(tf, ">> WHY NOT a single AI model for everything?")
body_line(tf, "Too slow for real-time. Too expensive at scale. We use AI periodically (v2), rules for real-time.")


# ── Slide 14: Detection Map ─────────────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "Detection Map: 6 Sub-Problems")
header_line(tf, ">> Each sub-problem, its technology, and update rate:")
body_line(tf, "1. FACIAL EXPRESSIONS   | Video frames   | MediaPipe FaceMesh (blendshapes)  | Every frame")
body_line(tf, "2. EYE CONTACT          | Video frames   | MediaPipe FaceMesh (gaze dirs)    | Every frame")
body_line(tf, "3. FILLER WORDS         | Transcript     | STT --> Pattern matching          | Every sentence")
body_line(tf, "4. HEDGING LANGUAGE     | Transcript     | Pattern matching (phrase list)    | Every sentence")
body_line(tf, "5. SPEAKING PACE        | Transcript+time| Word count / elapsed time         | Every 5 seconds")
body_line(tf, "6. VOICE STEADINESS     | Raw audio      | Web Audio API + FFT              | Continuous")
header_line(tf, ">> How they feed the score")
body_line(tf, "Detections 1+2 --> Face Score (x0.40)  |  Detections 3+4+5 --> Speech Score (x0.35)  |  Detection 6 --> Voice Score (x0.25)")
header_line(tf, ">> Optional v2 addition")
body_line(tf, "7. AI DEEP ANALYSIS     | Frame+text     | Claude API (multimodal)           | Every 30 seconds")
note_line(tf, "Detections 1 and 2 share the same MediaPipe model -- zero extra cost for eye contact.")


# ── Slide 15: 16 Blendshapes We Use ─────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "16 Blendshapes We Use (out of 52)")
header_line(tf, ">> Smile & Expression")
body_line(tf, "mouthSmileLeft / mouthSmileRight   (>0.08 = smiling)    |  mouthFrownLeft / mouthFrownRight  (>0.12 = frown)")
body_line(tf, "mouthPressLeft / mouthPressRight    (>0.08 = lip press)  |  jawOpen  (>0.02 = speaking)")
header_line(tf, ">> Brow & Worry")
body_line(tf, "browInnerUp  (>0.25 = worry/concern)  |  browDownLeft / browDownRight  (>0.30 = tension/frustration)")
header_line(tf, ">> Eyes: Blink & Gaze")
body_line(tf, "eyeBlinkLeft / eyeBlinkRight  (>0.50 = blink detected)  |  eyeWideLeft / eyeWideRight  (>0.15 = wide/surprised)")
body_line(tf, "eyeLookDownLeft/Right  (>0.55 = looking down)  |  eyeLookUpLeft/Right  (>0.55 = looking up)")
body_line(tf, "eyeLookInLeft/Right  (>0.55 = looking sideways)  |  eyeLookOutLeft/Right  (>0.55 = looking sideways)")
header_line(tf, ">> WHY only 16?")
body_line(tf, "The other 36 blendshapes cover tongue, cheek puff, lip pucker -- irrelevant to confidence.")
body_line(tf, "Fewer signals = less noise in scoring.")
note_line(tf, "Eye contact threshold is 0.55 (not 0.50) because laptop cameras sit slightly below eye level. 0.50 flagged normal laptop gaze as 'looking down'.")


# ── Slide 16: Filler & Hedge Lists ─────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "Filler Words & Hedging Phrases")
header_line(tf, ">> 18 Filler Words (5 categories)")
body_line(tf, "PURE VERBAL PAUSES:     um, uh, uhh, ah, er, hmm  (always fillers, zero meaning)")
body_line(tf, "CONTEXT-DEPENDENT:      like, so, well  (sometimes real words -- accept 15% false positive)")
body_line(tf, "VERBAL PADDING:         basically, actually, literally, honestly  (usually padding)")
body_line(tf, "FILLER PHRASES:         you know, I mean, okay so  (almost always fillers in presentations)")
body_line(tf, "SOFTENERS:              sort of, kind of, right  (soften statements, reduce authority)")
header_line(tf, ">> 16 Hedging Phrases (5 categories)")
body_line(tf, "UNCERTAINTY:            'I think', 'I believe', 'I feel like'  (opinion as uncertain)")
body_line(tf, "WEAK MODIFIERS:         'maybe', 'probably', 'perhaps'  (sounds unsure)")
body_line(tf, "DOWNPLAYING:            'sort of', 'kind of', 'a little bit'  (minimizing own statement)")
body_line(tf, "EXPLICIT DOUBT:         'I'm not sure', 'I could be wrong', 'I guess'  (stating lack of confidence)")
body_line(tf, "APOLOGIZING:            'sorry', 'sorry but'  (apologizing for having something to say)")
warn_line(tf, "Each filler: penalty = rate% x 5  |  Each hedge: -3 points  |  Caps prevent one signal from dominating")


# ── Slide 17: Comparison to Products ────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "Comparison to Existing Products")
header_line(tf, ">> Yoodli (AI Speech Coach)")
body_line(tf, "Cloud-based, requires internet, subscription pricing, closed source, good accuracy")
header_line(tf, ">> Poised (Real-time Feedback)")
body_line(tf, "Desktop app, works during Zoom calls, subscription, proprietary algorithms")
header_line(tf, ">> Orai (Mobile Speech Coach)")
body_line(tf, "Mobile-only, focuses on pace/energy/confidence, subscription model")
header_line(tf, ">> Our Confidence Detector")
green_line(tf, "OFFLINE: Runs 100% in the browser. No data leaves your device.")
green_line(tf, "OPEN SOURCE: Every algorithm is visible. Every threshold is tunable.")
green_line(tf, "FREE: No subscription. No API cost for core features.")
green_line(tf, "EDUCATIONAL: Built to learn AI/ML. Every decision is documented with WHY.")
green_line(tf, "CUSTOMIZABLE: Adjust weights, thresholds, filler lists to your needs.")
note_line(tf, "We trade polish for transparency. Yoodli is more polished -- but you can't see how it works or tune it.")


# ── Slide 18: Honest Limitations ────────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "Honest Limitations")
header_line(tf, ">> What this system CANNOT do:")
warn_line(tf, "1. NO SARCASM DETECTION: 'Oh great, another meeting' scores as positive (words are fine)")
warn_line(tf, "2. SINGLE SPEAKER ONLY: Cannot handle two people talking simultaneously")
warn_line(tf, "3. ACCENT-DEPENDENT: STT accuracy drops for non-American/British English accents")
warn_line(tf, "4. LIGHTING-DEPENDENT: MediaPipe needs decent lighting. Dark rooms = face not detected")
warn_line(tf, "5. 'LIKE' FALSE POSITIVES: 'I like pizza' penalized as filler (~15% false positive rate)")
warn_line(tf, "6. CULTURAL BIAS: 130-160 WPM 'optimal' is Western norm. Some cultures speak faster/slower")
warn_line(tf, "7. NO BODY LANGUAGE: Camera sees face only. Can't detect hand wringing below frame")
warn_line(tf, "8. MIC DEPENDENCY: Voice score varies with mic quality, distance, background noise")
header_line(tf, ">> Our approach to limitations")
body_line(tf, "Be transparent: tell users. Let users calibrate to their baseline. Never present")
body_line(tf, "scores as absolute truth -- always 'relative to your normal.'")
note_line(tf, "Acknowledging limitations is not weakness. It builds trust and sets correct expectations.")


# ── Slide 19: Learning Phases ───────────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "Learning Phases: 5-Phase Plan")
header_line(tf, ">> Phase 1: Hardware Access (Camera + Mic)")
body_line(tf, "getUserMedia API, permissions, video/audio stream lifecycle")
body_line(tf, "REUSED from ExamGuard: Camera access pattern, permission handling")
header_line(tf, ">> Phase 2: Speech & Audio")
body_line(tf, "Speech-to-Text (Vosk/Web Speech API), Web Audio API, pitch/volume extraction")
body_line(tf, "NEW: Everything. ExamGuard never touched audio.")
header_line(tf, ">> Phase 3: Face & Expression Detection")
body_line(tf, "MediaPipe FaceMesh, blendshapes, expression mapping, eye contact, blink rate")
body_line(tf, "REUSED from ExamGuard: MediaPipe setup, landmark reading, gaze detection")
header_line(tf, ">> Phase 4: NLP Text Analysis")
body_line(tf, "Tokenization, filler detection, hedging detection, pace calculation, text scoring")
body_line(tf, "NEW: Everything. ExamGuard never analyzed text.")
header_line(tf, ">> Phase 5: System Integration")
body_line(tf, "Multi-engine orchestration, scoring algorithm, dashboard, coaching alerts, reports")
body_line(tf, "REUSED from ExamGuard: Dashboard patterns, scoring concepts. NEW: multi-engine sync.")


# ── Slide 20: MVP Versions ──────────────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "MVP Versions: Build in Layers")
header_line(tf, ">> v0.1 -- Face Only (3-5 days)")
body_line(tf, "Camera --> MediaPipe Face Mesh --> Expression label + face overlay")
body_line(tf, "Proves: 'MediaPipe works in my browser and I can read facial expressions'")
header_line(tf, ">> v0.2 -- Face + Speech (1 week)")
body_line(tf, "Add microphone + STT + NLP. Two separate scores displayed side by side.")
body_line(tf, "Proves: 'Two engines can run simultaneously without lag'")
header_line(tf, ">> v0.3 -- Full Demo (2-3 weeks)")
body_line(tf, "Add voice engine + combined scoring + coaching alerts + session report")
body_line(tf, "Proves: 'This is a real product that coaches you in real time'")
header_line(tf, ">> v1.0 -- Production (2 weeks)")
body_line(tf, "Add Claude API feedback + session history + baseline calibration + PDF export + polish")
body_line(tf, "Proves: 'I built a portfolio-quality AI speech coaching product'")
warn_line(tf, "Each version is a WORKING product you can demo. Never build everything at once.")
note_line(tf, "v0.3 = 'It works.'   v1.0 = 'It works AND it looks like a real product.'")


# ── Slide 21: How to Get Accuracy ───────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "How to Get Accuracy")
header_line(tf, ">> The only path to a good score: test with real data")
body_line(tf, "Step 1: RECORD 5-10 test sessions with different confidence levels")
body_line(tf, "  - One session where you TRY to be confident")
body_line(tf, "  - One session with deliberate 'ums' everywhere")
body_line(tf, "  - One session avoiding eye contact")
body_line(tf, "  - One session at your normal comfort level")
body_line(tf, "Step 2: SCORE each session yourself (human judgment, 0-100). Be honest.")
body_line(tf, "Step 3: COMPARE human score vs system score")
body_line(tf, "  Example: 'Lots of ums' -- Human says 35, System says 55 --> system too generous")
body_line(tf, "Step 4: ADJUST weights to close the gap")
body_line(tf, "  If system is too generous on fillers --> increase Speech weight")
body_line(tf, "  If system ignores bad eye contact --> increase Face weight or adjust penalty")
body_line(tf, "Step 5: REPEAT until human and system agree within 5-10 points")
warn_line(tf, "No amount of theoretical weight-choosing beats testing with real data.")
note_line(tf, "This is the real path to accuracy. Plan 2 dedicated days for tuning after v0.3 is built.")


# ── Slide 22: Error Handling ────────────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "Error Handling: Degrade Gracefully")
header_line(tf, ">> Every failure scenario and what happens:")
body_line(tf, "CAMERA DENIED:       Show clear message. Continue with voice+speech only (reweight to 0.58/0.42)")
body_line(tf, "MIC DENIED:          Show clear message. Continue with face-only scoring")
body_line(tf, "BOTH DENIED:         Cannot score. Block session with helpful permission instructions.")
body_line(tf, "FACE NOT DETECTED:   Warn 'check lighting/angle'. Fall back to voice+speech scoring")
body_line(tf, "STT TIMEOUT (>10s):  Auto-restart STT silently. After 3 failures, warn user. Freeze speech score.")
body_line(tf, "NO AUDIO (>15s):     Prompt 'Are you still presenting?' Voice penalizes for silence. Face continues.")
body_line(tf, "TAB HIDDEN:          Pause session (camera stops when tab hidden). Resume on return.")
body_line(tf, "SLOW DEVICE:         Reduce to 15 FPS, then 10 FPS. Below 10: warn user. Accuracy reduced.")
header_line(tf, ">> Design principle")
green_line(tf, "Something is ALWAYS better than nothing. A voice-only score is still useful.")
warn_line(tf, "NEVER crash. NEVER show a blank screen. NEVER silently fail.")


# ── Slide 23: Performance Budget ────────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "Performance Budget")
header_line(tf, ">> Target: 30 FPS = 33.3ms per frame")
body_line(tf, "Camera sends 30 frames per second. Each frame must be processed in <34ms.")
body_line(tf, "But we don't need to process EVERY frame for face analysis.")
header_line(tf, ">> Strategy: Process every 2nd frame")
body_line(tf, "Process face on frames 1, 3, 5, 7...  Skip frames 2, 4, 6, 8...")
body_line(tf, "Effective rate: 15 FPS for face analysis. Still smooth, half the CPU cost.")
body_line(tf, "Audio engines run on their own thread (Web Audio API). No frame budget conflict.")
header_line(tf, ">> How to measure: Chrome DevTools")
body_line(tf, "Performance tab --> Record 10 seconds --> Look at frame times")
body_line(tf, "If frames consistently >34ms: optimize. If occasional spikes: acceptable.")
body_line(tf, "Use Web Workers for heavy processing to keep UI thread responsive.")
header_line(tf, ">> Fallback plan")
body_line(tf, "If 15 FPS not achievable: reduce resolution (640x480), process every 3rd frame,")
body_line(tf, "or drop voice engine first (lowest weight at 0.25).")
warn_line(tf, "Profile BEFORE adding features. Performance issues found late are 10x harder to fix.")


# ── Slide 24: What's Next ──────────────────────────────────────────
s = make_slide(prs)
tf = content_frame(s)
title_line(tf, "What's Next: The Build Order")
header_line(tf, ">> Step 1: Build v0.1 (Face Only)")
body_line(tf, "Camera + MediaPipe + expression labels. Prove the core tech works. 3-5 days.")
header_line(tf, ">> Step 2: Build v0.2 (Face + Speech)")
body_line(tf, "Add microphone + STT + NLP. Two engines running side by side. 1 week.")
header_line(tf, ">> Step 3: Build v0.3 (Full Demo)")
body_line(tf, "Add voice engine + combined scoring + coaching alerts + session reports. 2-3 weeks.")
header_line(tf, ">> Step 4: Build v1.0 (Production)")
body_line(tf, "Add Claude API + session history + calibration + PDF export + polish. 2 weeks.")
header_line(tf, ">> The philosophy")
green_line(tf, "A simple working product beats a complex unfinished product.", size=20)
body_line(tf, "Each version is a complete, demo-able product. Never build all 4 engines at once.")
body_line(tf, "When blocked: cut scope, never cut quality. A 2-engine system that scores accurately")
body_line(tf, "is better than a 4-engine system that gives wrong numbers.")
note_line(tf, "Total estimated build time: 6-8 weeks from Phase 1 to v1.0. You've done harder things in ExamGuard.")


# ── Save ────────────────────────────────────────────────────────────
output_path = r"D:\AI Learning\usman\07_Confidence_Detector_Project\07_Confidence_Detector_Project.pptx"
prs.save(output_path)
print(f"Saved presentation to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
